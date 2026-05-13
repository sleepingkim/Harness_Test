"""
3분 구두 발표 PPTX 생성기
- 영어 버전: oral_presentation_EN.pptx
- 한국어 버전: oral_presentation_KO.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ──────────────────────────────────────────
C_BLUE   = RGBColor(0x1F, 0x77, 0xB4)   # RAG / Improved
C_GRAY   = RGBColor(0x7F, 0x7F, 0x7F)   # Baseline / LLM-only
C_RED    = RGBColor(0xD6, 0x27, 0x28)   # 한계 / 문제
C_GREEN  = RGBColor(0x2C, 0xA0, 0x2C)   # 긍정 결과
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # 제목 텍스트
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)   # 슬라이드 배경
C_NAVY   = RGBColor(0x16, 0x32, 0x5B)   # 헤더 배경
C_ACCENT = RGBColor(0xE8, 0xF4, 0xFF)   # 하이라이트 배경

W = Inches(13.33)   # 와이드 16:9
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # 완전 빈 레이아웃
    return prs.slides.add_slide(layout)


def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h,
                size=18, bold=False, color=C_DARK,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle="", timer=""):
    """상단 네이비 헤더"""
    bar = add_rect(slide, 0, 0, W, Inches(1.2), fill=C_NAVY)
    add_textbox(slide, title,
                Inches(0.3), Inches(0.05), Inches(10), Inches(0.7),
                size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.3), Inches(0.72), Inches(9), Inches(0.4),
                    size=16, color=RGBColor(0xAD, 0xD8, 0xE6))
    if timer:
        add_textbox(slide, timer,
                    Inches(11.5), Inches(0.35), Inches(1.5), Inches(0.5),
                    size=14, color=RGBColor(0xAD, 0xD8, 0xE6),
                    align=PP_ALIGN.RIGHT)


def footer_bar(slide, author, event, slide_n, total):
    """하단 푸터"""
    bar = add_rect(slide, 0, Inches(7.1), W, Inches(0.4), fill=C_NAVY)
    txt = f"{author}  |  {event}  |  Slide {slide_n}/{total}"
    add_textbox(slide, txt,
                Inches(0.3), Inches(7.12), Inches(12), Inches(0.3),
                size=11, color=RGBColor(0xAD, 0xD8, 0xE6))


def highlight_box(slide, text, l, t, w, h,
                  bg=C_ACCENT, tc=C_BLUE, size=16, bold=True):
    add_rect(slide, l, t, w, h, fill=bg, line=C_BLUE)
    add_textbox(slide, text, l+Inches(0.1), t+Inches(0.05),
                w-Inches(0.2), h-Inches(0.1),
                size=size, bold=bold, color=tc, wrap=True)


def bullet_block(slide, items, l, t, w, h, size=16, color=C_DARK, gap=0.38):
    """items: list of (bullet_char, text, color)"""
    y = t
    for bc, txt, col in items:
        add_textbox(slide, bc, l, y, Inches(0.3), Inches(0.35),
                    size=size, bold=True, color=col)
        add_textbox(slide, txt, l+Inches(0.35), y, w-Inches(0.4), Inches(0.38),
                    size=size, color=col, wrap=True)
        y += Inches(gap)


# ══════════════════════════════════════════════════════════
#  슬라이드 콘텐츠 (EN / KO)
# ══════════════════════════════════════════════════════════

def build_slide1_en(slide):
    header_bar(slide,
               "The Problem: A Public Database Full of Aliases",
               "Same tool — dozens of names — one broken database",
               "⏱ 30s")
    # 클러스터 박스 3개
    clusters = [
        ("드릴비트세트", ["기리셋트", "giri set", "기리셑트", "drill bit SET"]),
        ("해머드릴",     ["함마드릴",  "함마 드릴", "HAMMER DRILL", "함마드릴기"]),
        ("파이프렌치",   ["PIPE렌치",  "파이프 렌치", "pipe wrench"]),
    ]
    cx = [Inches(0.4), Inches(4.7), Inches(9.0)]
    for i, (std, variants) in enumerate(clusters):
        bx = cx[i]
        # 표준명 박스
        std_box = add_rect(slide, bx, Inches(1.4), Inches(4.0), Inches(0.5), fill=C_BLUE)
        add_textbox(slide, std, bx+Inches(0.05), Inches(1.42), Inches(3.9), Inches(0.45),
                    size=15, bold=True, color=C_WHITE)
        # 비표준 표기들
        for j, v in enumerate(variants):
            vy = Inches(2.05) + j * Inches(0.55)
            add_rect(slide, bx+Inches(0.15), vy, Inches(3.7), Inches(0.42), fill=RGBColor(0xFF,0xE5,0xE5), line=C_RED)
            add_textbox(slide, v, bx+Inches(0.25), vy+Inches(0.04), Inches(3.5), Inches(0.38),
                        size=13, color=C_RED)

    # 통계 표
    add_rect(slide, Inches(0.4), Inches(4.5), Inches(5.5), Inches(0.45), fill=RGBColor(0xE8,0xF4,0xFF))
    add_textbox(slide, "Non-standard entries: 3,352   |   Standard names: 160   |   ~21 aliases per standard name",
                Inches(0.5), Inches(4.53), Inches(5.3), Inches(0.4), size=13, bold=True, color=C_BLUE)

    highlight_box(slide,
                  "⚠  Broken standardization → broken inventory, broken search, broken statistics.",
                  Inches(0.4), Inches(5.15), Inches(12.5), Inches(0.55),
                  bg=RGBColor(0xFF,0xF3,0xCD), tc=RGBColor(0x8B,0x60,0x00), size=14)

    footer_bar(slide, "Hongchul Shin", "UST 2026 Academic Workshop", 1, 6)


def build_slide1_ko(slide):
    header_bar(slide,
               "문제 정의: 비표준 표기로 가득 찬 공공 데이터베이스",
               "같은 공구 — 수십 가지 이름 — 하나의 손상된 데이터베이스",
               "⏱ 30초")
    clusters = [
        ("드릴비트세트", ["기리셋트", "giri set", "기리셑트", "drill bit SET"]),
        ("해머드릴",     ["함마드릴",  "함마 드릴", "HAMMER DRILL", "함마드릴기"]),
        ("파이프렌치",   ["PIPE렌치",  "파이프 렌치", "pipe wrench"]),
    ]
    cx = [Inches(0.4), Inches(4.7), Inches(9.0)]
    for i, (std, variants) in enumerate(clusters):
        bx = cx[i]
        add_rect(slide, bx, Inches(1.4), Inches(4.0), Inches(0.5), fill=C_BLUE)
        add_textbox(slide, std, bx+Inches(0.05), Inches(1.42), Inches(3.9), Inches(0.45),
                    size=15, bold=True, color=C_WHITE)
        for j, v in enumerate(variants):
            vy = Inches(2.05) + j * Inches(0.55)
            add_rect(slide, bx+Inches(0.15), vy, Inches(3.7), Inches(0.42), fill=RGBColor(0xFF,0xE5,0xE5), line=C_RED)
            add_textbox(slide, v, bx+Inches(0.25), vy+Inches(0.04), Inches(3.5), Inches(0.38),
                        size=13, color=C_RED)

    add_rect(slide, Inches(0.4), Inches(4.5), Inches(5.5), Inches(0.45), fill=RGBColor(0xE8,0xF4,0xFF))
    add_textbox(slide, "비표준 항목: 3,352개   |   표준 공구명: 160종   |   표준명당 약 21개 별칭",
                Inches(0.5), Inches(4.53), Inches(5.3), Inches(0.4), size=13, bold=True, color=C_BLUE)

    highlight_box(slide,
                  "⚠  비표준화 → 재고 집계 오류, 검색 실패, 통계 왜곡",
                  Inches(0.4), Inches(5.15), Inches(12.5), Inches(0.55),
                  bg=RGBColor(0xFF,0xF3,0xCD), tc=RGBColor(0x8B,0x60,0x00), size=14)

    footer_bar(slide, "신홍철 (Hongchul Shin)", "UST 2026 Academic Workshop", 1, 6)


# ── Slide 2 ──────────────────────────────────────────────

def build_slide2_en(slide):
    header_bar(slide,
               "Why Is This Hard? — Five Error Patterns",
               "78% of errors fall into 5 systematic patterns — none fixable by simple rules",
               "⏱ 25s")

    rows = [
        ("Typos",             "경랑몽키, 스판너",       33.9, C_DARK),
        ("Brand Inclusion",   "보쉬전동드릴",            15.5, C_DARK),
        ("Power Source",      "유선전동드릴",            10.9, C_DARK),
        ("English / Abbrev",  "HSS드릴비트",              9.0, C_DARK),
        ("Attribute Info",    "180mm그라인더",            8.8, C_DARK),
    ]
    bar_max = Inches(5.5)
    y0 = Inches(1.35)
    for name, ex, pct, col in rows:
        y = y0
        add_textbox(slide, name, Inches(0.4), y, Inches(2.5), Inches(0.38), size=14, color=col)
        bw = bar_max * pct / 34
        add_rect(slide, Inches(3.0), y+Inches(0.04), bw, Inches(0.3), fill=C_BLUE)
        add_textbox(slide, f"{pct}%", Inches(3.0)+bw+Inches(0.1), y, Inches(0.7), Inches(0.38), size=13, bold=True, color=C_BLUE)
        add_textbox(slide, ex, Inches(9.3), y, Inches(3.7), Inches(0.38), size=12, color=C_GRAY)
        y0 += Inches(0.52)

    # 합계 선
    add_rect(slide, Inches(0.4), Inches(4.05), Inches(12.5), Inches(0.02), fill=C_DARK)
    add_textbox(slide, "Top 5 Combined: 78.1%", Inches(0.4), Inches(4.1), Inches(4), Inches(0.38),
                size=14, bold=True, color=C_DARK)

    # 일본어 잔재 강조
    add_rect(slide, Inches(0.4), Inches(4.55), Inches(12.5), Inches(0.6), fill=RGBColor(0xFF,0xE5,0xE5), line=C_RED)
    add_textbox(slide, "⚠  Japanese-derived terms  (4.6%):  기리 / 뺀치 / 함마",
                Inches(0.55), Inches(4.6), Inches(7), Inches(0.35), size=14, bold=True, color=C_RED)
    add_textbox(slide, "Phonetic discontinuity → automation impossible without a thesaurus",
                Inches(0.55), Inches(4.95), Inches(9), Inches(0.3), size=13, color=C_RED)

    footer_bar(slide, "Hongchul Shin", "UST 2026 Academic Workshop", 2, 6)


def build_slide2_ko(slide):
    header_bar(slide,
               "왜 어려운가? — 5가지 오류 유형",
               "오류의 78%가 5가지 체계적 패턴 — 단순 규칙으로는 해결 불가",
               "⏱ 25초")

    rows = [
        ("오탈자",          "경랑몽키, 스판너",      33.9),
        ("브랜드 혼입",     "보쉬전동드릴",           15.5),
        ("동력원 명시",     "유선전동드릴",           10.9),
        ("영문·약어",       "HSS드릴비트",             9.0),
        ("속성 기입",       "180mm그라인더",           8.8),
    ]
    bar_max = Inches(5.5)
    y0 = Inches(1.35)
    for name, ex, pct in rows:
        y = y0
        add_textbox(slide, name, Inches(0.4), y, Inches(2.5), Inches(0.38), size=14, color=C_DARK)
        bw = bar_max * pct / 34
        add_rect(slide, Inches(3.0), y+Inches(0.04), bw, Inches(0.3), fill=C_BLUE)
        add_textbox(slide, f"{pct}%", Inches(3.0)+bw+Inches(0.1), y, Inches(0.7), Inches(0.38), size=13, bold=True, color=C_BLUE)
        add_textbox(slide, ex, Inches(9.3), y, Inches(3.7), Inches(0.38), size=12, color=C_GRAY)
        y0 += Inches(0.52)

    add_rect(slide, Inches(0.4), Inches(4.05), Inches(12.5), Inches(0.02), fill=C_DARK)
    add_textbox(slide, "상위 5개 유형 합계: 78.1%", Inches(0.4), Inches(4.1), Inches(4), Inches(0.38),
                size=14, bold=True, color=C_DARK)

    add_rect(slide, Inches(0.4), Inches(4.55), Inches(12.5), Inches(0.6), fill=RGBColor(0xFF,0xE5,0xE5), line=C_RED)
    add_textbox(slide, "⚠  일본어 잔재 (4.6%):  기리 / 뺀치 / 함마",
                Inches(0.55), Inches(4.6), Inches(7), Inches(0.35), size=14, bold=True, color=C_RED)
    add_textbox(slide, "음운 단절 → 동의어 사전 없이 자동화 불가",
                Inches(0.55), Inches(4.95), Inches(9), Inches(0.3), size=13, color=C_RED)

    footer_bar(slide, "신홍철 (Hongchul Shin)", "UST 2026 Academic Workshop", 2, 6)


# ── Slide 3 ──────────────────────────────────────────────

def build_slide3_en(slide):
    header_bar(slide,
               "Our Approach: RAG v2 Pipeline",
               "Three targeted upgrades over the baseline",
               "⏱ 35s")

    # 파이프라인 박스들
    boxes = [
        (Inches(0.3),  "Raw Input\n(tool name)",    C_GRAY),
        (Inches(3.4),  "Retrieval\n(Top-k=15)",      C_BLUE),
        (Inches(6.5),  "LLM Judgment\n(local model)", C_BLUE),
        (Inches(9.6),  "Standard Name\nOutput",       C_GREEN),
    ]
    bw, bh = Inches(2.8), Inches(0.85)
    by = Inches(1.45)
    for bx, label, col in boxes:
        add_rect(slide, bx, by, bw, bh, fill=col)
        add_textbox(slide, label, bx+Inches(0.1), by+Inches(0.08), bw-Inches(0.2), bh-Inches(0.15),
                    size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 화살표 (텍스트로)
    for ax in [Inches(3.15), Inches(6.25), Inches(9.35)]:
        add_textbox(slide, "►", ax, by+Inches(0.22), Inches(0.3), Inches(0.45),
                    size=20, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    # Fallback
    add_textbox(slide, "LLM fails?  ▼  Fallback: return RAG rank-1",
                Inches(6.5), Inches(2.42), Inches(5.8), Inches(0.38),
                size=12, color=C_BLUE)

    # 업그레이드 비교표
    tbl_y = Inches(3.0)
    headers = ["Component", "Baseline (gray)", "RAG v2 Improved (blue)"]
    col_w = [Inches(2.5), Inches(4.5), Inches(5.8)]
    col_x = [Inches(0.3), Inches(2.9), Inches(7.5)]

    # 헤더
    for i, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        c = C_NAVY if i == 0 else (C_GRAY if i == 1 else C_BLUE)
        add_rect(slide, cx, tbl_y, cw, Inches(0.38), fill=c)
        add_textbox(slide, hdr, cx+Inches(0.05), tbl_y+Inches(0.04), cw-Inches(0.1), Inches(0.32),
                    size=13, bold=True, color=C_WHITE)

    rows = [
        ("Embedding",  "MiniLM  (384-dim)",          "BGE-m3-ko  (1024-dim)  ★"),
        ("Retrieval",  "FAISS only,  k=5",            "BM25 40% + FAISS 60%,  k=15"),
        ("Safety net", "—",                           "Fallback (rank-1 auto-return)"),
    ]
    for ri, (comp, base, imp) in enumerate(rows):
        ry = tbl_y + Inches(0.38) + ri * Inches(0.46)
        bg = C_LGRAY if ri % 2 == 0 else C_WHITE
        for ci, (txt, cx, cw) in enumerate(zip([comp, base, imp], col_x, col_w)):
            add_rect(slide, cx, ry, cw, Inches(0.44), fill=bg)
            col = C_DARK if ci < 2 else C_BLUE
            bld = ci == 2
            add_textbox(slide, txt, cx+Inches(0.08), ry+Inches(0.05), cw-Inches(0.15), Inches(0.36),
                        size=13, bold=bld, color=col)

    footer_bar(slide, "Hongchul Shin", "UST 2026 Academic Workshop", 3, 6)


def build_slide3_ko(slide):
    header_bar(slide,
               "제안 방법: RAG v2 파이프라인",
               "베이스라인 대비 세 가지 핵심 업그레이드",
               "⏱ 35초")

    boxes = [
        (Inches(0.3),  "입력\n(공구명)",              C_GRAY),
        (Inches(3.4),  "검색\n(상위 k=15 후보)",       C_BLUE),
        (Inches(6.5),  "LLM 판단\n(로컬 모델)",        C_BLUE),
        (Inches(9.6),  "표준 공구명\n출력",             C_GREEN),
    ]
    bw, bh = Inches(2.8), Inches(0.85)
    by = Inches(1.45)
    for bx, label, col in boxes:
        add_rect(slide, bx, by, bw, bh, fill=col)
        add_textbox(slide, label, bx+Inches(0.1), by+Inches(0.08), bw-Inches(0.2), bh-Inches(0.15),
                    size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    for ax in [Inches(3.15), Inches(6.25), Inches(9.35)]:
        add_textbox(slide, "►", ax, by+Inches(0.22), Inches(0.3), Inches(0.45),
                    size=20, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    add_textbox(slide, "LLM 실패 시  ▼  Fallback: RAG 1위 후보 자동 반환",
                Inches(6.5), Inches(2.42), Inches(6.0), Inches(0.38),
                size=12, color=C_BLUE)

    tbl_y = Inches(3.0)
    headers = ["구성 요소", "베이스라인", "RAG v2 개선 버전"]
    col_w = [Inches(2.5), Inches(4.5), Inches(5.8)]
    col_x = [Inches(0.3), Inches(2.9), Inches(7.5)]

    for i, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        c = C_NAVY if i == 0 else (C_GRAY if i == 1 else C_BLUE)
        add_rect(slide, cx, tbl_y, cw, Inches(0.38), fill=c)
        add_textbox(slide, hdr, cx+Inches(0.05), tbl_y+Inches(0.04), cw-Inches(0.1), Inches(0.32),
                    size=13, bold=True, color=C_WHITE)

    rows = [
        ("임베딩",    "MiniLM (384차원)",           "BGE-m3-ko (1024차원)  ★"),
        ("검색",      "FAISS 단독,  k=5",           "BM25 40% + FAISS 60%,  k=15"),
        ("안전장치",  "—",                          "Fallback (1위 자동 반환)"),
    ]
    for ri, (comp, base, imp) in enumerate(rows):
        ry = tbl_y + Inches(0.38) + ri * Inches(0.46)
        bg = C_LGRAY if ri % 2 == 0 else C_WHITE
        for ci, (txt, cx, cw) in enumerate(zip([comp, base, imp], col_x, col_w)):
            add_rect(slide, cx, ry, cw, Inches(0.44), fill=bg)
            col = C_DARK if ci < 2 else C_BLUE
            add_textbox(slide, txt, cx+Inches(0.08), ry+Inches(0.05), cw-Inches(0.15), Inches(0.36),
                        size=13, bold=(ci == 2), color=col)

    footer_bar(slide, "신홍철 (Hongchul Shin)", "UST 2026 Academic Workshop", 3, 6)


# ── Slide 4 ──────────────────────────────────────────────

def build_slide4_en(slide):
    header_bar(slide,
               "The Verdict: RAG Transforms Accuracy",
               "RAG is not optional — it is the difference between 9% and 88%",
               "⏱ 40s")

    models = [
        ("gemma4:e4b",      27.4, 88.47, "+61.1pp", True),
        ("gemma3:4b",       23.0, 84.48, "+61.5pp", False),
        ("exaone3.5:7.8b",   9.4, 79.16, "+69.8pp", False),
    ]
    bar_scale = Inches(4.2) / 100
    y0 = Inches(1.4)
    for model, llm_pct, rag_pct, delta, is_best in models:
        # 모델명
        add_textbox(slide, model, Inches(0.3), y0, Inches(2.2), Inches(0.35),
                    size=13, bold=is_best, color=C_DARK)
        # LLM-only 바
        lw = bar_scale * llm_pct
        add_rect(slide, Inches(2.6), y0+Inches(0.04), lw, Inches(0.28), fill=C_GRAY)
        add_textbox(slide, f"LLM-only {llm_pct}%",
                    Inches(2.6), y0+Inches(0.32), Inches(2), Inches(0.28),
                    size=11, color=C_GRAY)
        # RAG 바
        rw = bar_scale * rag_pct
        rc = C_GREEN if is_best else C_BLUE
        add_rect(slide, Inches(2.6), y0+Inches(0.65), rw, Inches(0.28), fill=rc)
        star = " ★ BEST" if is_best else ""
        add_textbox(slide, f"RAG  {rag_pct}%{star}",
                    Inches(2.6), y0+Inches(0.93), Inches(2.5), Inches(0.28),
                    size=11, bold=is_best, color=rc)
        # 브래킷
        bx = Inches(2.6) + rw + Inches(0.1)
        add_textbox(slide, delta,
                    bx, y0+Inches(0.3), Inches(1.2), Inches(0.6),
                    size=13, bold=True, color=C_BLUE)
        y0 += Inches(1.4)

    # DeepSeek 역설
    highlight_box(slide,
                  "DeepSeek Paradox:  Korean generation  0%  (LLM-only)  →  78.05%  with RAG\n"
                  "→  RAG compensates for absent language capability.",
                  Inches(0.3), Inches(5.7), Inches(7.8), Inches(0.78),
                  bg=C_ACCENT, tc=C_BLUE, size=13)

    add_textbox(slide, "* 451 test samples · Exact Match Accuracy",
                Inches(9.0), Inches(6.7), Inches(4.0), Inches(0.3),
                size=10, color=C_GRAY)

    footer_bar(slide, "Hongchul Shin", "UST 2026 Academic Workshop", 4, 6)


def build_slide4_ko(slide):
    header_bar(slide,
               "결과: RAG가 정확도를 바꾼다",
               "RAG는 선택이 아니라 필수 — 9%와 88%의 차이",
               "⏱ 40초")

    models = [
        ("gemma4:e4b",      27.4, 88.47, "+61.1pp", True),
        ("gemma3:4b",       23.0, 84.48, "+61.5pp", False),
        ("exaone3.5:7.8b",   9.4, 79.16, "+69.8pp", False),
    ]
    bar_scale = Inches(4.2) / 100
    y0 = Inches(1.4)
    for model, llm_pct, rag_pct, delta, is_best in models:
        add_textbox(slide, model, Inches(0.3), y0, Inches(2.2), Inches(0.35),
                    size=13, bold=is_best, color=C_DARK)
        lw = bar_scale * llm_pct
        add_rect(slide, Inches(2.6), y0+Inches(0.04), lw, Inches(0.28), fill=C_GRAY)
        add_textbox(slide, f"LLM 단독 {llm_pct}%",
                    Inches(2.6), y0+Inches(0.32), Inches(2), Inches(0.28),
                    size=11, color=C_GRAY)
        rw = bar_scale * rag_pct
        rc = C_GREEN if is_best else C_BLUE
        add_rect(slide, Inches(2.6), y0+Inches(0.65), rw, Inches(0.28), fill=rc)
        star = " ★ 최우수" if is_best else ""
        add_textbox(slide, f"RAG  {rag_pct}%{star}",
                    Inches(2.6), y0+Inches(0.93), Inches(2.5), Inches(0.28),
                    size=11, bold=is_best, color=rc)
        bx = Inches(2.6) + rw + Inches(0.1)
        add_textbox(slide, delta, bx, y0+Inches(0.3), Inches(1.2), Inches(0.6),
                    size=13, bold=True, color=C_BLUE)
        y0 += Inches(1.4)

    highlight_box(slide,
                  "DeepSeek 역설:  한국어 생성 능력  0%  (LLM 단독)  →  RAG 적용 후  78.05%\n"
                  "→  RAG가 언어 능력 자체의 한계를 우회한다.",
                  Inches(0.3), Inches(5.7), Inches(7.8), Inches(0.78),
                  bg=C_ACCENT, tc=C_BLUE, size=13)

    add_textbox(slide, "* 451개 테스트 샘플 · Exact Match Accuracy",
                Inches(9.0), Inches(6.7), Inches(4.0), Inches(0.3),
                size=10, color=C_GRAY)

    footer_bar(slide, "신홍철 (Hongchul Shin)", "UST 2026 Academic Workshop", 4, 6)


# ── Slide 5 ──────────────────────────────────────────────

def build_slide5_en(slide):
    header_bar(slide,
               "What Drives Performance? — Embedding Matters Most",
               "Same LLM. Different pipeline. +22~27 percentage points.",
               "⏱ 30s")

    # 비교표
    tbl_y = Inches(1.45)
    for ci, (label, col) in enumerate([("Model", C_NAVY),
                                        ("Baseline  (MiniLM + FAISS k=5)", C_GRAY),
                                        ("RAG v2 Improved", C_BLUE)]):
        cw = [Inches(2.5), Inches(4.5), Inches(5.8)][ci]
        cx = [Inches(0.3), Inches(2.9), Inches(7.5)][ci]
        add_rect(slide, cx, tbl_y, cw, Inches(0.38), fill=col)
        add_textbox(slide, label, cx+Inches(0.05), tbl_y+Inches(0.04), cw-Inches(0.1), Inches(0.32),
                    size=13, bold=True, color=C_WHITE)

    rows_cmp = [
        ("granite4.1:8b", "55.21%", "82.48%  (+27.3pp)"),
        ("deepseek-r1:8b","55.21%", "77.83%  (+22.6pp)"),
    ]
    for ri, (m, base, imp) in enumerate(rows_cmp):
        ry = tbl_y + Inches(0.38) + ri * Inches(0.46)
        bg = C_LGRAY if ri == 0 else C_WHITE
        for ci, (txt, cx, cw) in enumerate(zip([m, base, imp],
                                                [Inches(0.3), Inches(2.9), Inches(7.5)],
                                                [Inches(2.5), Inches(4.5), Inches(5.8)])):
            add_rect(slide, cx, ry, cw, Inches(0.44), fill=bg)
            col = C_DARK if ci < 2 else C_BLUE
            add_textbox(slide, txt, cx+Inches(0.08), ry+Inches(0.05), cw-Inches(0.15), Inches(0.36),
                        size=14, bold=(ci == 2), color=col)

    # 요인 분해
    add_textbox(slide, "Performance gain decomposition (estimated):",
                Inches(0.3), Inches(3.6), Inches(7), Inches(0.35), size=14, bold=True, color=C_DARK)

    factors = [
        ("Embedding swap  (MiniLM → BGE-m3-ko)", "~15–20pp", 20, C_BLUE, True),
        ("Candidate pool  (k=5 → k=15)",          "~3–5pp",    5, C_NAVY, False),
        ("BM25 hybrid  (FAISS → BM25+FAISS)",     "~3–5pp",    5, C_NAVY, False),
    ]
    scale = Inches(6.0) / 25
    y0 = Inches(4.05)
    for label, val, bar_v, col, dominant in factors:
        bw = scale * bar_v
        add_rect(slide, Inches(0.3), y0+Inches(0.05), bw, Inches(0.32), fill=col)
        add_textbox(slide, label, Inches(0.3)+bw+Inches(0.1), y0, Inches(5.5), Inches(0.38),
                    size=13, color=col)
        add_textbox(slide, val, Inches(7.5), y0, Inches(1.5), Inches(0.38),
                    size=13, bold=dominant, color=col)
        if dominant:
            add_textbox(slide, "← DOMINANT", Inches(9.1), y0, Inches(2), Inches(0.38),
                        size=12, bold=True, color=C_BLUE)
        y0 += Inches(0.5)

    highlight_box(slide,
                  "Embedding selection dominates pipeline performance.\nChoose your embedding before tuning your LLM.",
                  Inches(0.3), Inches(5.65), Inches(12.5), Inches(0.72),
                  bg=C_ACCENT, tc=C_BLUE, size=14)

    footer_bar(slide, "Hongchul Shin", "UST 2026 Academic Workshop", 5, 6)


def build_slide5_ko(slide):
    header_bar(slide,
               "무엇이 성능을 결정하는가? — 임베딩이 핵심",
               "같은 LLM, 다른 파이프라인 — +22~27 퍼센트포인트",
               "⏱ 30초")

    tbl_y = Inches(1.45)
    for ci, (label, col) in enumerate([("모델", C_NAVY),
                                        ("베이스라인  (MiniLM + FAISS k=5)", C_GRAY),
                                        ("RAG v2 개선 버전", C_BLUE)]):
        cw = [Inches(2.5), Inches(4.5), Inches(5.8)][ci]
        cx = [Inches(0.3), Inches(2.9), Inches(7.5)][ci]
        add_rect(slide, cx, tbl_y, cw, Inches(0.38), fill=col)
        add_textbox(slide, label, cx+Inches(0.05), tbl_y+Inches(0.04), cw-Inches(0.1), Inches(0.32),
                    size=13, bold=True, color=C_WHITE)

    rows_cmp = [
        ("granite4.1:8b", "55.21%", "82.48%  (+27.3pp)"),
        ("deepseek-r1:8b","55.21%", "77.83%  (+22.6pp)"),
    ]
    for ri, (m, base, imp) in enumerate(rows_cmp):
        ry = tbl_y + Inches(0.38) + ri * Inches(0.46)
        bg = C_LGRAY if ri == 0 else C_WHITE
        for ci, (txt, cx, cw) in enumerate(zip([m, base, imp],
                                                [Inches(0.3), Inches(2.9), Inches(7.5)],
                                                [Inches(2.5), Inches(4.5), Inches(5.8)])):
            add_rect(slide, cx, ry, cw, Inches(0.44), fill=bg)
            col = C_DARK if ci < 2 else C_BLUE
            add_textbox(slide, txt, cx+Inches(0.08), ry+Inches(0.05), cw-Inches(0.15), Inches(0.36),
                        size=14, bold=(ci == 2), color=col)

    add_textbox(slide, "성능 향상 요인 분해 (추정):",
                Inches(0.3), Inches(3.6), Inches(7), Inches(0.35), size=14, bold=True, color=C_DARK)

    factors = [
        ("임베딩 교체  (MiniLM → BGE-m3-ko)", "~15–20pp", 20, C_BLUE, True),
        ("후보 확대  (k=5 → k=15)",            "~3–5pp",    5, C_NAVY, False),
        ("BM25 하이브리드 추가",               "~3–5pp",    5, C_NAVY, False),
    ]
    scale = Inches(6.0) / 25
    y0 = Inches(4.05)
    for label, val, bar_v, col, dominant in factors:
        bw = scale * bar_v
        add_rect(slide, Inches(0.3), y0+Inches(0.05), bw, Inches(0.32), fill=col)
        add_textbox(slide, label, Inches(0.3)+bw+Inches(0.1), y0, Inches(5.5), Inches(0.38),
                    size=13, color=col)
        add_textbox(slide, val, Inches(7.5), y0, Inches(1.5), Inches(0.38),
                    size=13, bold=dominant, color=col)
        if dominant:
            add_textbox(slide, "← 지배적", Inches(9.1), y0, Inches(2), Inches(0.38),
                        size=12, bold=True, color=C_BLUE)
        y0 += Inches(0.5)

    highlight_box(slide,
                  "임베딩 선택이 파이프라인 성능을 지배한다.\nLLM 튜닝보다 임베딩 선택이 먼저다.",
                  Inches(0.3), Inches(5.65), Inches(12.5), Inches(0.72),
                  bg=C_ACCENT, tc=C_BLUE, size=14)

    footer_bar(slide, "신홍철 (Hongchul Shin)", "UST 2026 Academic Workshop", 5, 6)


# ── Slide 6 ──────────────────────────────────────────────

def build_slide6_en(slide):
    header_bar(slide,
               "Conclusion & Takeaways",
               "Three things to remember.",
               "⏱ 20s")

    bullets = [
        ("①", "RAG is essential",      "+61–70pp accuracy gain over LLM-only",           C_BLUE),
        ("②", "Small model wins",      "4B gemma4:e4b → 88.47% exact match",             C_GREEN),
        ("③", "Embedding first",       "BGE-m3-ko swap alone → +22–27pp",               C_BLUE),
    ]
    y0 = Inches(1.5)
    for num, title, desc, col in bullets:
        add_rect(slide, Inches(0.3), y0, Inches(0.55), Inches(0.55), fill=col)
        add_textbox(slide, num, Inches(0.3), y0+Inches(0.06), Inches(0.55), Inches(0.45),
                    size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, Inches(0.95), y0+Inches(0.04), Inches(3.5), Inches(0.38),
                    size=17, bold=True, color=col)
        add_textbox(slide, desc, Inches(4.6), y0+Inches(0.08), Inches(8.3), Inches(0.38),
                    size=15, color=C_DARK)
        y0 += Inches(0.82)

    # 한계
    add_rect(slide, Inches(0.3), Inches(4.05), Inches(12.5), Inches(0.6),
             fill=RGBColor(0xFF,0xE5,0xE5), line=C_RED)
    add_textbox(slide,
                "⚠  Remaining challenge:  Japanese-derived terms → 47.1%  |  Synonym dictionary integration needed",
                Inches(0.45), Inches(4.12), Inches(12.1), Inches(0.45),
                size=13, color=C_RED)

    highlight_box(slide,
                  '"Privacy-safe, cost-effective, and fully deployable on-premise."',
                  Inches(0.3), Inches(5.0), Inches(12.5), Inches(0.75),
                  bg=RGBColor(0xE8,0xF5,0xE9), tc=C_GREEN, size=18)

    footer_bar(slide, "Hongchul Shin", "UST 2026 Academic Workshop", 6, 6)


def build_slide6_ko(slide):
    header_bar(slide,
               "결론 및 핵심 정리",
               "기억할 세 가지.",
               "⏱ 20초")

    bullets = [
        ("①", "RAG는 필수",       "LLM 단독 대비 정확도 +61~70pp 향상",              C_BLUE),
        ("②", "소형 모델로 충분", "4B gemma4:e4b → Exact Match 88.47%",             C_GREEN),
        ("③", "임베딩이 먼저",    "BGE-m3-ko 교체만으로 +22~27pp",                   C_BLUE),
    ]
    y0 = Inches(1.5)
    for num, title, desc, col in bullets:
        add_rect(slide, Inches(0.3), y0, Inches(0.55), Inches(0.55), fill=col)
        add_textbox(slide, num, Inches(0.3), y0+Inches(0.06), Inches(0.55), Inches(0.45),
                    size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, Inches(0.95), y0+Inches(0.04), Inches(3.5), Inches(0.38),
                    size=17, bold=True, color=col)
        add_textbox(slide, desc, Inches(4.6), y0+Inches(0.08), Inches(8.3), Inches(0.38),
                    size=15, color=C_DARK)
        y0 += Inches(0.82)

    add_rect(slide, Inches(0.3), Inches(4.05), Inches(12.5), Inches(0.6),
             fill=RGBColor(0xFF,0xE5,0xE5), line=C_RED)
    add_textbox(slide,
                "⚠  남은 과제:  일본어 잔재 → 47.1% 오류율  |  동의어 사전 통합 필요",
                Inches(0.45), Inches(4.12), Inches(12.1), Inches(0.45),
                size=13, color=C_RED)

    highlight_box(slide,
                  '"프라이버시 안전, 비용 효율적, 온프레미스 완전 배포 가능."',
                  Inches(0.3), Inches(5.0), Inches(12.5), Inches(0.75),
                  bg=RGBColor(0xE8,0xF5,0xE9), tc=C_GREEN, size=18)

    footer_bar(slide, "신홍철 (Hongchul Shin)", "UST 2026 Academic Workshop", 6, 6)


# ══════════════════════════════════════════════════════════
#  타이틀 슬라이드 (Slide 0)
# ══════════════════════════════════════════════════════════

def build_title_en(slide):
    add_rect(slide, 0, 0, W, H, fill=C_NAVY)
    add_rect(slide, 0, Inches(4.8), W, Inches(0.06), fill=C_BLUE)

    add_textbox(slide,
                "Semantic Standardization of Unstructured\nCategorical Data Using Local LLMs and RAG",
                Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.9),
                size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "A Case Study on Industrial Tool Names",
                Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.6),
                size=22, color=RGBColor(0xAD,0xD8,0xE6), align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "Hongchul Shin\nUST-KIST HDI Lab",
                Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.8),
                size=18, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "KIST School, UST 2026 Academic Workshop  |  3-Minute Oral Presentation",
                Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.5),
                size=14, color=RGBColor(0xAD,0xD8,0xE6), align=PP_ALIGN.CENTER)


def build_title_ko(slide):
    add_rect(slide, 0, 0, W, H, fill=C_NAVY)
    add_rect(slide, 0, Inches(4.8), W, Inches(0.06), fill=C_BLUE)

    add_textbox(slide,
                "비정형 범주형 데이터의 의미적 표준화\n로컬 LLM과 RAG를 활용하여",
                Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.9),
                size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "서울시 공구대여 데이터 공구명 표준화 사례 연구",
                Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.6),
                size=20, color=RGBColor(0xAD,0xD8,0xE6), align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "신홍철 (Hongchul Shin)\nUST-KIST HDI 연구실",
                Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.8),
                size=18, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "KIST School, UST 2026 Academic Workshop  |  3분 구두 발표",
                Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.5),
                size=14, color=RGBColor(0xAD,0xD8,0xE6), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
#  메인
# ══════════════════════════════════════════════════════════

def build(lang):
    prs = new_prs()

    builders_en = [build_title_en, build_slide1_en, build_slide2_en,
                   build_slide3_en, build_slide4_en, build_slide5_en, build_slide6_en]
    builders_ko = [build_title_ko, build_slide1_ko, build_slide2_ko,
                   build_slide3_ko, build_slide4_ko, build_slide5_ko, build_slide6_ko]

    builders = builders_en if lang == "EN" else builders_ko

    for fn in builders:
        slide = blank_slide(prs)
        fn(slide)

    out = f"C:/Users/neohc/Desktop/ClaudeCode/_workspace5/oral_presentation/oral_presentation_{lang}.pptx"
    prs.save(out)
    print(f"[OK] {lang} saved -> {out}")


if __name__ == "__main__":
    build("EN")
    build("KO")
    print("Done.")
