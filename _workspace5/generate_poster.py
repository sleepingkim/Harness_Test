#!/usr/bin/env python3
"""
Academic Poster Generator — LLM+RAG Semantic Standardization
Output: poster_draft.pptx  (48" x 36" landscape, A0-equivalent)
Run  : python generate_poster.py
Needs: pip install python-pptx
"""

import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── COLOR PALETTE ─────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x3A, 0x6B)
MID_BLUE    = RGBColor(0x2E, 0x70, 0xBF)
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2A)
PANEL_BG    = RGBColor(0xFD, 0xFD, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF2, 0xF5)
WARN_BG     = RGBColor(0xFF, 0xF3, 0xCD)
WARN_TEXT   = RGBColor(0x7D, 0x42, 0x00)
ALT_ROW     = RGBColor(0xE8, 0xF1, 0xFB)

# ─── POSTER DIMENSIONS (48" × 36" landscape) ──────────────────────────────────
W      = Inches(48)
H      = Inches(36)
MARGIN = Inches(0.4)
GAP    = Inches(0.3)


# ─── LOW-LEVEL HELPERS ─────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, bg, border=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)   # 1 = RECTANGLE
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(0.75)
    return shp


def textbox(slide, l, t, w, h, text,
            size=24, bold=False, italic=False,
            color=WHITE, align=PP_ALIGN.LEFT,
            wrap=True, name="Calibri"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text          = text
    run.font.size     = Pt(size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = color
    run.font.name     = name
    return txb


def multitext(slide, l, t, w, h, lines,
              size=20, color=DARK_TEXT, name="Calibri", spacing=Pt(2)):
    """Multiple-paragraph textbox; each item in `lines` becomes a paragraph."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = spacing
        run = p.add_run()
        run.text           = line
        run.font.size      = Pt(size)
        run.font.color.rgb = color
        run.font.name      = name
    return txb


def section_hdr(slide, l, t, w, label, size=30):
    """Returns bottom-y of the header bar."""
    hh = Inches(0.65)
    rect(slide, l, t, w, hh, MID_BLUE)
    textbox(slide, l + Inches(0.12), t + Inches(0.08),
            w - Inches(0.24), hh - Inches(0.1),
            label, size=size, bold=True, color=WHITE, wrap=False)
    return t + hh


def sub_hdr(slide, l, t, w, label, size=22):
    """Thin accent bar for sub-sections. Returns bottom-y."""
    hh = Inches(0.45)
    rect(slide, l, t, w, hh, LIGHT_BLUE)
    textbox(slide, l + Inches(0.1), t + Inches(0.05),
            w - Inches(0.2), hh - Inches(0.08),
            label, size=size, bold=True, color=DARK_BLUE, wrap=False)
    return t + hh


def styled_table(slide, l, t, w, data, font=17, hdr_font=19):
    """Draw a styled table from a list-of-lists. Returns bottom-y."""
    rows = len(data)
    cols = len(data[0])
    row_h = Inches(0.42)
    tbl   = slide.shapes.add_table(rows, cols, l, t, w, row_h * rows).table

    # equal column widths by default
    col_w = int(w) // cols
    for c in range(cols):
        tbl.columns[c].width = col_w
    for r in range(rows):
        tbl.rows[r].height = row_h

    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            tf  = cell.text_frame
            p   = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = "Calibri"
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
                run.font.bold = True; run.font.size = Pt(hdr_font)
                run.font.color.rgb = WHITE
            else:
                bg = ALT_ROW if r % 2 == 0 else WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
                run.font.size = Pt(font)
                run.font.color.rgb = DARK_TEXT

    return t + row_h * rows


# ─── POSTER CONSTRUCTION ───────────────────────────────────────────────────────

def build(out_path):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    # ── BACKGROUND ─────────────────────────────────────────────────────────────
    rect(slide, 0, 0, W, H, LIGHT_GRAY)

    # ══════════════════════════════════════════════════════════════════════
    # HEADER BANNER
    # ══════════════════════════════════════════════════════════════════════
    HDR_H = Inches(4.4)
    rect(slide, 0, 0, W, HDR_H, DARK_BLUE)

    # Title
    textbox(slide,
        l=Inches(0.6), t=Inches(0.25), w=Inches(46.8), h=Inches(2.3),
        text="Semantic Standardization of Unstructured Categorical Data\n"
             "Using Local LLMs and Retrieval-Augmented Generation",
        size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    textbox(slide,
        l=Inches(0.6), t=Inches(2.6), w=Inches(46.8), h=Inches(0.75),
        text="A Case Study on Industrial Tool Name Standardization "
             "with Open-Source LLMs",
        size=32, color=RGBColor(0xAD, 0xD8, 0xF7),
        align=PP_ALIGN.CENTER)

    # Authors
    textbox(slide,
        l=Inches(0.6), t=Inches(3.4), w=Inches(46.8), h=Inches(0.7),
        text="Hongcheol Shin   |   Korea University of Science and Technology (UST)   "
             "|   ETRI (Electronics and Telecommunications Research Institute)",
        size=25, color=RGBColor(0xCF, 0xE2, 0xFF), align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════
    # ABSTRACT  (full width)
    # ══════════════════════════════════════════════════════════════════════
    AY = HDR_H
    AH = Inches(3.6)
    rect(slide, 0, AY, W, AH, WHITE)
    ay = section_hdr(slide, MARGIN, AY + Inches(0.15),
                     W - 2*MARGIN, "ABSTRACT", size=32)
    multitext(slide,
        l=MARGIN + Inches(0.2), t=ay + Inches(0.1),
        w=W - 2*MARGIN - Inches(0.4), h=AH - Inches(1.1),
        lines=[
            "Unstructured categorical data from manual entry exhibits severe lexical variability, undermining data quality and downstream machine learning. "
            "This work proposes a Retrieval-Augmented Generation (RAG) pipeline employing open-source LLMs to semantically standardize "
            "industrial tool names from the Seoul Metropolitan Tool Rental Dataset (3,352 unique dirty strings → 160 standard categories). "
            "We evaluate 12 open-source LLMs (1.5B–9B) under a baseline RAG (MiniLM-L12, FAISS, k=5) and an improved RAG "
            "(BGE-m3-ko embedding, BM25+FAISS hybrid retrieval, k=15, category-augmented prompting). "
            "The improved configuration reaches up to 88.47% accuracy (+33.3 pp over baseline ≈55%). "
            "We also identify and correct a train/test data leakage flaw in prior work (RAG v1: 97.88% was inflated). "
            "Results show that domain-optimized retrieval architecture matters more than raw model scale."
        ], size=21, color=DARK_TEXT)

    # ══════════════════════════════════════════════════════════════════════
    # THREE-COLUMN BODY
    # ══════════════════════════════════════════════════════════════════════
    CY  = AY + AH + Inches(0.15)
    CH  = Inches(18.8)
    CW  = (W - 2*MARGIN - 2*GAP) / 3

    C1  = MARGIN
    C2  = MARGIN + CW + GAP
    C3  = MARGIN + 2*(CW + GAP)

    # Column panel backgrounds
    for cx in (C1, C2, C3):
        rect(slide, cx, CY, CW, CH, PANEL_BG,
             border=RGBColor(0xC0, 0xC8, 0xD8))

    # ── COLUMN 1 : INTRODUCTION ───────────────────────────────────────────
    PAD = Inches(0.18)
    y1  = CY + Inches(0.12)
    y1  = section_hdr(slide, C1+PAD, y1, CW-2*PAD, "1. INTRODUCTION", size=28)
    y1 += Inches(0.12)

    multitext(slide,
        l=C1+PAD+Inches(0.08), t=y1, w=CW-2*PAD-Inches(0.16), h=Inches(2.5),
        lines=[
            "Background",
            "Real-world categorical fields entered manually suffer from severe lexical variability. "
            "Seoul's public tool rental database contains 3,352 unique tool name strings that collapse "
            "to only ~160 real categories (≈21 variants per category on average).",
            "",
            "Limitation of Rule-Based Approaches",
            "Regex / morphological analyzers address syntactic errors only. They cannot handle semantic "
            "variability: Japanese-derived legacy terms (giri → drill bit), embedded brand/spec attributes, "
            "compound nouns, or context-dependent disambiguation.",
        ], size=19, color=DARK_TEXT)
    y1 += Inches(2.65)

    y1 = sub_hdr(slide, C1+PAD, y1, CW-2*PAD, "Dirty Data Taxonomy (13 sub-types)", size=21)
    y1 += Inches(0.08)
    dirty = [
        ["Type",              "Example",               "N"],
        ["Attribute-embedded","'180mm Sander'",         "946"],
        ["Brand inclusion",   "'BOSCH push drive'",     "209"],
        ["Power source noted","'Cordless drill 20V'",   "330"],
        ["Japanese-derived",  "'giri' → drill bit",     "87"],
        ["Typo / encoding",   "'gyeongrang monkey'",    "117"],
        ["Set / compound",    "'glue gun & refill'",    "315"],
        ["Metadata / noise",  "'(1) hand saw'",         "147"],
    ]
    y1 = styled_table(slide, C1+PAD, y1, CW-2*PAD, dirty, font=16, hdr_font=17)
    y1 += Inches(0.18)

    y1 = sub_hdr(slide, C1+PAD, y1, CW-2*PAD, "Research Contributions", size=21)
    y1 += Inches(0.12)
    multitext(slide,
        l=C1+PAD+Inches(0.08), t=y1, w=CW-2*PAD-Inches(0.16), h=Inches(3.5),
        lines=[
            "① Rigorous evaluation framework",
            "   Strict separation of retrieval knowledge base from ground-truth evaluation set, "
            "   eliminating the train/test data leakage present in prior work (RAG v1).",
            "",
            "② Improved RAG pipeline",
            "   Four targeted enhancements: Korean-specialized embedding, BM25+FAISS hybrid "
            "   retrieval, expanded candidate pool (k=15), and category-augmented prompting.",
            "",
            "③ Comparative LLM benchmark",
            "   12 open-source LLMs (1.5B–9B) evaluated under both baseline and improved "
            "   configurations on a real-world Korean dirty-data corpus.",
        ], size=18, color=DARK_TEXT)

    # ── COLUMN 2 : MATERIALS & METHODS ──────────────────────────────────
    y2  = CY + Inches(0.12)
    y2  = section_hdr(slide, C2+PAD, y2, CW-2*PAD, "2. MATERIALS & METHODS", size=28)
    y2 += Inches(0.12)

    y2 = sub_hdr(slide, C2+PAD, y2, CW-2*PAD, "2.1 Dataset", size=21)
    y2 += Inches(0.08)
    ds = [
        ["Item",                    "Value"],
        ["Source",                  "Seoul Tool Rental Open Dataset"],
        ["Raw records",             "11,525 rental entries"],
        ["Unique dirty strings",    "3,352"],
        ["Standard categories",     "160"],
        ["Avg. variants / category","~21"],
        ["Ground-truth pairs",      "457 (manually curated)"],
    ]
    y2 = styled_table(slide, C2+PAD, y2, CW-2*PAD, ds, font=17, hdr_font=18)
    y2 += Inches(0.18)

    y2 = sub_hdr(slide, C2+PAD, y2, CW-2*PAD, "2.2 RAG Architecture", size=21)
    y2 += Inches(0.08)
    cfg = [
        ["Component",    "Baseline",                  "Improved (★)"],
        ["Embedding",    "MiniLM-L12  (384-dim)",     "BGE-m3-ko  (1024-dim)  ★"],
        ["Retrieval",    "FAISS only",                "BM25 (40%) + FAISS (60%)  ★"],
        ["Candidates k", "5",                         "15  ★"],
        ["Prompt",       "Standard",                  "+ Category injection  ★"],
        ["Knowledge DB", "160 standard names",        "160 standard names"],
        ["Eval set",     "457 GT pairs (separated)",  "457 GT pairs (separated)"],
    ]
    y2 = styled_table(slide, C2+PAD, y2, CW-2*PAD, cfg, font=17, hdr_font=18)
    y2 += Inches(0.1)

    multitext(slide,
        l=C2+PAD+Inches(0.08), t=y2, w=CW-2*PAD-Inches(0.16), h=Inches(1.0),
        lines=[
            "Hybrid score = 0.4 × BM25_norm + 0.6 × FAISS_cosine_norm",
            "BM25 tokenizer: Korean syllable unigram + bigram  "
            "(e.g., '경랑몽키' → ['경','랑','몽','키','경랑','랑몽','몽키'])",
        ], size=17, color=DARK_TEXT)
    y2 += Inches(1.1)

    y2 = sub_hdr(slide, C2+PAD, y2, CW-2*PAD, "2.3 Models Evaluated  (12 total)", size=21)
    y2 += Inches(0.08)
    mdls = [
        ["Family",   "Models",                              "Params"],
        ["Gemma",    "gemma3:4b, gemma4:e2b, gemma4:e4b",  "4B / 2B / 4B"],
        ["Granite",  "granite4.1:3b, granite4.1:8b",       "3B / 8B"],
        ["EXAONE",   "exaone3.5:2.4b, exaone3.5:7.8b",     "2.4B / 7.8B"],
        ["DeepSeek", "deepseek-r1:1.5b, deepseek-r1:8b",   "1.5B / 8B"],
        ["Qwen",     "qwen3.5:9b",                         "9B"],
    ]
    y2 = styled_table(slide, C2+PAD, y2, CW-2*PAD, mdls, font=17, hdr_font=18)
    y2 += Inches(0.16)

    multitext(slide,
        l=C2+PAD+Inches(0.08), t=y2, w=CW-2*PAD-Inches(0.16), h=Inches(1.5),
        lines=[
            "Hardware:  NVIDIA RTX 5080 (17.1 GB VRAM)  |  Ollama (local GPU inference)",
            "Embedding: GPU-accelerated batch processing (3,352 queries → ~10 sec)",
            "Metrics:   Accuracy (exact match)  |  Fallback rate  |  Speed (items/sec)",
        ], size=18, color=DARK_TEXT)

    # ── COLUMN 3 : RESULTS ────────────────────────────────────────────────
    y3  = CY + Inches(0.12)
    y3  = section_hdr(slide, C3+PAD, y3, CW-2*PAD, "3. RESULTS", size=28)
    y3 += Inches(0.12)

    y3 = sub_hdr(slide, C3+PAD, y3, CW-2*PAD, "3.1 Overall Performance (12 models)", size=21)
    y3 += Inches(0.08)
    res = [
        ["Rank","Model",              "Config",  "Accuracy","Fallback","Speed"],
        ["1","gemma4:e4b",          "improved","88.47%",  "8.1%",   "1.58/s"],
        ["2","gemma3:4b",           "improved","84.48%",  "6.5%",   "1.77/s"],
        ["3","granite4.1:8b",       "improved","82.48%",  "5.0%",   "1.21/s"],
        ["4","gemma4:e2b",          "improved","80.93%",  "6.4%",   "1.93/s"],
        ["5","exaone3.5:7.8b",      "improved","79.16%",  "9.2%",   "1.23/s"],
        ["6","deepseek-r1:1.5b",    "improved","78.05%",  "85.2%*", "1.49/s"],
        ["7","deepseek-r1:8b",      "improved","77.83%",  "2.6%",   "1.40/s"],
        ["8","exaone3.5:2.4b",      "improved","74.06%",  "13.6%",  "1.58/s"],
        ["9","granite4.1:3b",       "improved","67.85%",  "20.8%",  "1.53/s"],
        ["10","deepseek-r1:8b",     "baseline","55.21%",  "3.8%",   "1.96/s"],
        ["11","granite4.1:8b",      "baseline","55.21%",  "14.9%",  "1.26/s"],
        ["12","qwen3.5:9b",         "baseline","0.00%",   "0.0%",   "1.31/s"],
    ]
    y3 = styled_table(slide, C3+PAD, y3, CW-2*PAD, res, font=15, hdr_font=16)
    y3 += Inches(0.1)

    # Warning note — v1 leakage
    rect(slide, C3+PAD, y3, CW-2*PAD, Inches(0.65), WARN_BG)
    multitext(slide,
        l=C3+PAD+Inches(0.12), t=y3+Inches(0.05),
        w=CW-2*PAD-Inches(0.24), h=Inches(0.55),
        lines=["⚠  RAG v1 (prior work) reported 97.88% — inflated by data leakage "
               "(eval pairs directly in retrieval DB). RAG v2 eliminates this flaw."],
        size=16, color=WARN_TEXT)
    y3 += Inches(0.75)

    y3 = sub_hdr(slide, C3+PAD, y3, CW-2*PAD, "3.2 Key Findings", size=21)
    y3 += Inches(0.1)
    multitext(slide,
        l=C3+PAD+Inches(0.08), t=y3, w=CW-2*PAD-Inches(0.16), h=Inches(5.5),
        lines=[
            "① Model size does NOT predict RAG performance",
            "   gemma3:4b (3.3 GB) > exaone3.5:7.8b (4.8 GB)  →  84.5% vs. 79.2%",
            "   deepseek-r1:1.5b (1.1 GB) achieves 78% — comparable to 8B models",
            "",
            "② All four improvements contribute",
            "   BGE-m3-ko: captures Korean morphology "
            "('경랑몽키' → 'monkey wrench')",
            "   BM25 hybrid: recovers lexical signals missed by pure vector search",
            "   k=15: granite4.1:8b fallback  14.9% → 5.0%",
            "   Category injection: improves disambiguation of near-synonym tools",
            "",
            "③ Fallback rate as instruction-following proxy",
            "   deepseek-r1:1.5b: 85.2% → format compliance failure",
            "   Top-3 models: 5–8% → acceptable for production deployment",
            "",
            "④ Speed — all top-9 models: < 50 min / 3,352 items on RTX 5080",
            "   (vs. 1–13 h on Mac Mini M1 CPU in prior work)",
        ], size=17, color=DARK_TEXT)

    # ══════════════════════════════════════════════════════════════════════
    # CONCLUSION + REFERENCES  (bottom row)
    # ══════════════════════════════════════════════════════════════════════
    BOT_Y = CY + CH + Inches(0.15)
    BOT_H = H - BOT_Y - Inches(0.1)

    CONC_W = W * 0.63
    REF_W  = W - MARGIN - CONC_W - MARGIN - GAP

    rect(slide, MARGIN, BOT_Y, CONC_W, BOT_H, PANEL_BG,
         border=RGBColor(0xC0, 0xC8, 0xD8))
    rect(slide, MARGIN + CONC_W + GAP, BOT_Y, REF_W, BOT_H, PANEL_BG,
         border=RGBColor(0xC0, 0xC8, 0xD8))

    # Conclusion
    cy = section_hdr(slide, MARGIN+PAD, BOT_Y+Inches(0.1),
                     CONC_W-2*PAD, "4. CONCLUSION", size=28)
    multitext(slide,
        l=MARGIN+PAD+Inches(0.1), t=cy+Inches(0.1),
        w=CONC_W-2*PAD-Inches(0.2), h=BOT_H-Inches(1.2),
        lines=[
            "Summary",
            "A RAG pipeline with four targeted improvements raises accuracy from ~55% (baseline) to 88.47% (gemma4:e4b), "
            "a +33.3 pp gain on a rigorous, leakage-free benchmark. Domain-optimized retrieval architecture outweighs raw "
            "model scale — compact models (3–4 GB) match or exceed 8B counterparts in this task. Local open-source LLMs "
            "provide a privacy-safe, cost-effective alternative to cloud APIs for domain-specific categorical data standardization.",
            "",
            "Limitations",
            "  · Ground-truth coverage: 457 / 3,352 = 13.6%  (evaluation may not fully represent the full dirty-data distribution)",
            "  · Models >17B parameters: impractical on a single RTX 5080 (CPU offload reduces speed to ~0.03 items/s)",
            "  · deepseek-r1:1.5b: 85.2% fallback rate limits practical applicability despite 78% accuracy",
            "",
            "Future Work",
            "  ① Large-model evaluation (32B+) with quantization or multi-GPU inference",
            "  ② Full ground-truth expansion to all 3,352 records",
            "  ③ Human-in-the-Loop for low-confidence predictions (ambiguous compounds, set items)",
            "  ④ Fine-tuning vs. RAG comparison on domain-specific tool name data",
            "  ⑤ Generalization to other dirty categorical domains "
            "(ICD medical codes, e-commerce SKUs, HR job titles)",
        ], size=19, color=DARK_TEXT)

    # References
    REF_X = MARGIN + CONC_W + GAP
    ry = section_hdr(slide, REF_X+PAD, BOT_Y+Inches(0.1),
                     REF_W-2*PAD, "REFERENCES", size=26)
    refs = [
        "[1] Zhang & Huang (2024). Data Cleaning Using LLMs. arXiv:2410.15547",
        "[2] Glock & Korom (2025). Detecting Errors in Contact Info with LLMs. VLDB WS",
        "[3] Amugongo et al. (2025). RAG for LLMs in Healthcare. PLOS Digital Health",
        "[4] Ovadia et al. (2023). Fine-Tuning or Retrieval? arXiv:2312.05934",
        "[5] Hwang et al. (2025). Domain-Specific Embedding for Korean. arXiv:2502.07131",
        "[6] Belcak et al. (2025). Small LMs are Future of Agentic AI. arXiv:2506.02153",
        "[7] Kim et al. (2003). A Taxonomy of Dirty Data. Data Mining & KD, 7, 81–99",
        "[8] Wong et al. (2025). PolyNorm: Few-Shot Text Normalization. arXiv:2511.03080",
    ]
    multitext(slide,
        l=REF_X+PAD+Inches(0.1), t=ry+Inches(0.1),
        w=REF_W-2*PAD-Inches(0.2), h=BOT_H-Inches(1.2),
        lines=refs, size=17, color=DARK_TEXT)

    # ── SAVE ─────────────────────────────────────────────────────────────
    prs.save(out_path)
    print(f"\n  poster saved: {out_path}\n")


# ─── ENTRY POINT ───────────────────────────────────────────────────────────────
if __name__ == "__main__":
    try:
        from pptx import Presentation        # noqa: F401  (import check)
    except ImportError:
        sys.exit("ERROR: python-pptx not found.\n"
                 "       Run:  pip install python-pptx")

    script_dir = os.path.dirname(os.path.abspath(__file__))
    out = os.path.join(script_dir, "poster_draft.pptx")
    build(out)
