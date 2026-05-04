"""
비정형 데이터 정제, LLM이 해결할 수 있는가?
발표 초안 v3 → PowerPoint 자동 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Cm
import copy

# ── 색상 팔레트 ──────────────────────────────
C_NAVY    = RGBColor(0x1B, 0x2A, 0x4A)   # 제목 배경 (진남색)
C_ACCENT  = RGBColor(0x2E, 0x86, 0xC1)   # 강조 (파랑)
C_ORANGE  = RGBColor(0xE6, 0x7E, 0x22)   # 2차 강조 (주황)
C_GREEN   = RGBColor(0x1E, 0x8B, 0x5B)   # 성공/개선 (초록)
C_RED     = RGBColor(0xC0, 0x39, 0x2B)   # 경고/한계 (빨강)
C_LGRAY   = RGBColor(0xF4, 0xF6, 0xF8)   # 슬라이드 배경 (연회색)
C_DKGRAY  = RGBColor(0x5D, 0x6D, 0x7E)   # 본문 보조 (중회색)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK   = RGBColor(0x17, 0x20, 0x2A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # 완전 빈 레이아웃


# ── 공통 헬퍼 ─────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=0):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = Pt(line_width)
    else:
        line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, color=C_BLACK,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def title_bar(slide, title, subtitle=None):
    """상단 제목 바 (진남색 배경)"""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.3), fill_rgb=C_NAVY)
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.12), Inches(12), Inches(0.75),
                 font_size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.82), Inches(12), Inches(0.42),
                     font_size=14, color=RGBColor(0xAE, 0xC6, 0xCF), italic=True)

def section_label(slide, text, l, t, w=Inches(2.5), h=Inches(0.35)):
    """섹션 레이블 (파랑 배경 + 흰 글자)"""
    add_rect(slide, l, t, w, h, fill_rgb=C_ACCENT)
    add_text_box(slide, text, l+Inches(0.08), t+Inches(0.03),
                 w-Inches(0.16), h-Inches(0.06),
                 font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def key_number(slide, number, label, l, t, w=Inches(2.2), h=Inches(1.5),
               num_color=C_ACCENT):
    """큰 숫자 + 설명 카드"""
    add_rect(slide, l, t, w, h, fill_rgb=C_WHITE,
             line_rgb=C_ACCENT, line_width=1.5)
    add_text_box(slide, number,
                 l+Inches(0.1), t+Inches(0.05), w-Inches(0.2), Inches(0.9),
                 font_size=36, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    add_text_box(slide, label,
                 l+Inches(0.1), t+Inches(0.92), w-Inches(0.2), Inches(0.52),
                 font_size=11, color=C_DKGRAY, align=PP_ALIGN.CENTER, wrap=True)

def bullet_box(slide, items, l, t, w, h, title=None, title_color=C_ACCENT,
               font_size=14, bullet="▪ "):
    """불릿 텍스트 박스"""
    cur_t = t
    if title:
        add_text_box(slide, title, l, cur_t, w, Inches(0.38),
                     font_size=15, bold=True, color=title_color)
        cur_t += Inches(0.38)
    for item in items:
        add_text_box(slide, bullet + item,
                     l + Inches(0.1), cur_t, w - Inches(0.1), Inches(0.36),
                     font_size=font_size, color=C_BLACK)
        cur_t += Inches(0.34)


# ════════════════════════════════════════════════════
# SLIDE 1: 표지
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_NAVY)

# 장식 라인
add_rect(sl, Inches(0), Inches(4.6), SLIDE_W, Inches(0.06), fill_rgb=C_ACCENT)
add_rect(sl, Inches(0), Inches(4.66), SLIDE_W, Inches(0.03), fill_rgb=C_ORANGE)

add_text_box(sl, "비정형 데이터 정제, LLM이 해결할 수 있는가?",
             Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.4),
             font_size=38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(sl, "— 로컬 LLM·RAG 기반 의미론적 표준화: 실험 결과·오류 분석·실무 적용 방향",
             Inches(0.8), Inches(2.85), Inches(11.7), Inches(0.6),
             font_size=18, color=RGBColor(0xAE, 0xC6, 0xCF),
             align=PP_ALIGN.CENTER, italic=True)

add_text_box(sl, "신홍철  |  2026.05",
             Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.5),
             font_size=16, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(sl, "기반: 전문석사 프렉티컴보고서 (RAG v1) + RAG v2 재실험",
             Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.45),
             font_size=13, color=C_DKGRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 2: 목차
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LGRAY)
title_bar(sl, "목차")

sections = [
    ("01", "데이터 정제는 왜 병목인가",          "더티데이터의 현실과 유형"),
    ("02", "LLM+RAG로 어떻게 해결하는가",        "파이프라인 설계 + v1 한계 극복"),
    ("03", "왜 로컬 LLM이어야 하는가",           "보안·비용·안정성·커스터마이징"),
    ("04", "속도 vs 정확도 트레이드오프",         "9개 모델 실험 결과 + 시나리오 분석"),
    ("05", "LLM은 어디서 실수하는가",            "842건 오류 유형화 + 원인 분석"),
    ("06", "실제 업무에 적용한다면",              "구현 로드맵 + 기대 효과"),
    ("07", "결론 및 제언",                       "핵심 메시지 + 향후 연구"),
]

for i, (num, title, sub) in enumerate(sections):
    row = i % 4
    col = i // 4
    l = Inches(0.6 + col * 6.4)
    t = Inches(1.55 + row * 1.38)
    add_rect(sl, l, t, Inches(5.8), Inches(1.15),
             fill_rgb=C_WHITE, line_rgb=C_ACCENT, line_width=1)
    add_rect(sl, l, t, Inches(0.65), Inches(1.15), fill_rgb=C_ACCENT)
    add_text_box(sl, num, l+Inches(0.02), t+Inches(0.3), Inches(0.61), Inches(0.55),
                 font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, title, l+Inches(0.75), t+Inches(0.1), Inches(5.0), Inches(0.42),
                 font_size=15, bold=True, color=C_NAVY)
    add_text_box(sl, sub,   l+Inches(0.75), t+Inches(0.56), Inches(5.0), Inches(0.42),
                 font_size=11, color=C_DKGRAY, italic=True)


# ════════════════════════════════════════════════════
# SLIDE 3: 데이터 정제 병목 — 현황 수치
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LGRAY)
title_bar(sl, "01. 데이터 정제는 왜 병목인가", "서울시 공구대여 데이터 사례")

# 좌측: 핵심 문장
add_rect(sl, Inches(0.4), Inches(1.5), Inches(5.6), Inches(1.1),
         fill_rgb=C_NAVY)
add_text_box(sl, '"Data Scientists spend 60~80% of their time\ncleaning data."',
             Inches(0.5), Inches(1.55), Inches(5.4), Inches(0.95),
             font_size=15, bold=True, color=C_WHITE, wrap=True)
add_text_box(sl, "출처: CrowdFlower Survey (2016) / Zhang & Huang (2024) [12]",
             Inches(0.5), Inches(2.62), Inches(5.6), Inches(0.32),
             font_size=10, color=C_DKGRAY, italic=True)

add_text_box(sl, "AI 프로젝트 실패 원인 1위:",
             Inches(0.5), Inches(3.1), Inches(5.6), Inches(0.38),
             font_size=15, bold=False, color=C_DKGRAY)
add_rect(sl, Inches(0.4), Inches(3.48), Inches(5.6), Inches(0.55), fill_rgb=C_RED)
add_text_box(sl, "\"모델이 나빠서\"가 아니라 \"입력 데이터가 더러워서\"",
             Inches(0.5), Inches(3.5), Inches(5.4), Inches(0.5),
             font_size=14, bold=True, color=C_WHITE)
add_text_box(sl, "정선우 외(2023): 더티데이터가 LLM 출력 신뢰도 직접 저하 [14]",
             Inches(0.5), Inches(4.1), Inches(5.6), Inches(0.32),
             font_size=10, color=C_DKGRAY, italic=True)

# 우측: 수치 카드 4개
add_text_box(sl, "서울시 공구대여 데이터",
             Inches(6.5), Inches(1.5), Inches(6.3), Inches(0.4),
             font_size=16, bold=True, color=C_NAVY)

cards = [
    ("11,525건",  "전체 대여 기록"),
    ("3,352개",   "고유 공구 이름"),
    ("~160종",    "실제 공구 종류 (표준명)"),
    ("~21가지",   "1종당 평균 표기 변형"),
]
for i, (num, lbl) in enumerate(cards):
    col = i % 2
    row = i // 2
    l = Inches(6.5 + col * 3.1)
    t = Inches(2.0 + row * 1.65)
    key_number(sl, num, lbl, l, t, w=Inches(2.8), h=Inches(1.45),
               num_color=C_ACCENT if i < 3 else C_ORANGE)

add_rect(sl, Inches(6.4), Inches(6.9), Inches(6.5), Inches(0.05), fill_rgb=C_ACCENT)


# ════════════════════════════════════════════════════
# SLIDE 4: 더티데이터 5대 유형
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LGRAY)
title_bar(sl, "01. 비정형 더티데이터의 유형 분류", "Kim et al. (2003) 분류 체계 공구 도메인 적용 [18]")

types = [
    ("속성·규격 혼입",  C_ACCENT,  "'180mm 그라인더'",     "크기·규격을 공구명에 포함"),
    ("동력원 명시",     C_GREEN,   "'충전식전동드릴'",      "충전식/유선/에어 등을 이름에 기재"),
    ("오탈자·오표기",   C_ORANGE,  "'경랑몽키' → 경량몽키", "입력 오류 또는 구어체 표기"),
    ("일본어 잔재",     C_RED,     "'기리' '뺀치' '노기스'", "건설 현장 일본어 유래 표현"),
    ("세트·복합",       C_NAVY,    "'글루건&글루스틱'",      "여러 공구를 하나의 이름으로"),
]

for i, (name, color, ex, desc) in enumerate(types):
    l = Inches(0.4 + i * 2.52)
    t = Inches(1.55)
    add_rect(sl, l, t, Inches(2.3), Inches(0.55), fill_rgb=color)
    add_text_box(sl, name, l+Inches(0.08), t+Inches(0.1),
                 Inches(2.14), Inches(0.38),
                 font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, l, t+Inches(0.55), Inches(2.3), Inches(3.5),
             fill_rgb=C_WHITE, line_rgb=color, line_width=1.5)
    add_text_box(sl, "예시:", l+Inches(0.12), t+Inches(0.7),
                 Inches(2.1), Inches(0.3), font_size=11, color=C_DKGRAY, bold=True)
    add_text_box(sl, ex, l+Inches(0.12), t+Inches(1.0),
                 Inches(2.1), Inches(0.7),
                 font_size=13, bold=True, color=color, wrap=True)
    add_text_box(sl, desc, l+Inches(0.12), t+Inches(1.75),
                 Inches(2.1), Inches(1.0),
                 font_size=11, color=C_DKGRAY, wrap=True)

add_text_box(sl, "→ 수작업 규칙화 불가: 3,352개 변형에 수백 개 규칙 필요 (Glock & Korom, 2025) [10]",
             Inches(0.4), Inches(5.6), Inches(12.5), Inches(0.4),
             font_size=13, color=C_RED, bold=True)


# ════════════════════════════════════════════════════
# SLIDE 5: LLM+RAG 파이프라인
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LGRAY)
title_bar(sl, "02. LLM+RAG 파이프라인", "비정형 공구명 → 160종 표준명 분류 문제로 재정의")

# 파이프라인 흐름
steps = [
    ("더티 공구명 입력",        "'충전식전동드릴 (보쉬 18V)'",   C_NAVY),
    ("임베딩 (GPU 배치)",       "dragonkue/BGE-m3-ko\n~11초/3,352건",  C_ACCENT),
    ("하이브리드 검색",          "FAISS(60%) + BM25(40%)\n후보 15개 생성", C_GREEN),
    ("LLM 분류",                "Ollama 로컬 추론\nJSON 출력",     C_ORANGE),
    ("표준화 결과 출력",         "standard_name + brand\n+ power_source + spec", C_NAVY),
]

arrow_w = Inches(0.35)
box_w   = Inches(2.2)
gap     = Inches(0.25)
start_l = Inches(0.35)
t_box   = Inches(1.6)

for i, (title, desc, col) in enumerate(steps):
    l = start_l + i * (box_w + arrow_w + gap)
    add_rect(sl, l, t_box, box_w, Inches(2.3), fill_rgb=col)
    add_text_box(sl, title, l+Inches(0.1), t_box+Inches(0.12),
                 box_w-Inches(0.2), Inches(0.55),
                 font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=True)
    add_text_box(sl, desc, l+Inches(0.1), t_box+Inches(0.75),
                 box_w-Inches(0.2), Inches(1.35),
                 font_size=11, color=RGBColor(0xD5, 0xE8, 0xF5),
                 align=PP_ALIGN.CENTER, wrap=True)
    if i < len(steps)-1:
        al = l + box_w + gap*0.1
        add_text_box(sl, "▶", al, t_box+Inches(0.85),
                     arrow_w, Inches(0.5),
                     font_size=22, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# v1 vs v2 비교
add_text_box(sl, "RAG v1 vs v2 핵심 차이",
             Inches(0.4), Inches(4.3), Inches(12.5), Inches(0.38),
             font_size=15, bold=True, color=C_NAVY)

compare = [
    ("벡터 DB 입력", "386개 정답 쌍 (누수)", "160개 표준명 목록만"),
    ("평가 데이터",  "동일한 386개 (중복)",  "457개 — DB와 완전 분리"),
    ("최고 정확도",  "97.88% (과대 추정)",   "88.47% (신뢰 가능)"),
]
hdrs = ["항목", "RAG v1 (문제)", "RAG v2 (개선)"]
col_w = [Inches(2.5), Inches(4.2), Inches(4.2)]
col_l = [Inches(0.5), Inches(3.1), Inches(7.4)]

for j, (hdr, cw, cl) in enumerate(zip(hdrs, col_w, col_l)):
    add_rect(sl, cl, Inches(4.7), cw, Inches(0.35),
             fill_rgb=C_NAVY if j==0 else (C_RED if j==1 else C_GREEN))
    add_text_box(sl, hdr, cl+Inches(0.05), Inches(4.72), cw-Inches(0.1), Inches(0.32),
                 font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

for r, (item, v1, v2) in enumerate(compare):
    row_t = Inches(5.05) + r * Inches(0.52)
    bg = C_LGRAY if r % 2 == 0 else C_WHITE
    for j, (val, cw, cl) in enumerate(zip([item, v1, v2], col_w, col_l)):
        add_rect(sl, cl, row_t, cw, Inches(0.48), fill_rgb=bg)
        col_txt = C_BLACK if j == 0 else (C_RED if j == 1 else C_GREEN)
        add_text_box(sl, val, cl+Inches(0.05), row_t+Inches(0.06),
                     cw-Inches(0.1), Inches(0.38),
                     font_size=11, color=col_txt, bold=(j > 0), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 6: 왜 로컬 LLM인가
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LGRAY)
title_bar(sl, "03. 왜 로컬 LLM이어야 하는가", "4가지 핵심 이유")

reasons = [
    ("🔒 보안·프라이버시",
     ["내부 데이터가 외부 서버로 전송되지 않음",
      "개인정보보호법·정보보안 정책 준수 가능",
      "Amugongo et al. (2025): 의료 RAG에서 로컬 배포 권고 [5]",
      "Belcak et al. (2025): 강력한 데이터 통제 실현 [7]"],
     C_NAVY),
    ("💰 비용 구조",
     ["3,352건: 클라우드 ~$7 → 로컬 ~$0.05",
      "월 100만 건: 클라우드 ~$2,100 → 로컬 ~$15",
      "WINGPT (2025): 규모의 경제 — 로컬 비용 우위 급증 [2]"],
     C_GREEN),
    ("⚡ 운영 안정성",
     ["인터넷 장애 시에도 정상 운영",
      "API Rate Limit 없음, 자유로운 배치 처리",
      "Belcak et al. (2025): 오프라인 에이전트 추론 시연 [7]"],
     C_ORANGE),
    ("🔧 커스터마이징",
     ["공구 표준 개정 시 RAG DB만 갱신",
      "도메인 지식 변경 시 모델 재학습 불필요",
      "Ovadia et al. (2023): 업데이트 잦은 도메인에 RAG 우월 [16]"],
     C_ACCENT),
]

for i, (title, bullets, color) in enumerate(reasons):
    col = i % 2
    row = i // 2
    l = Inches(0.4 + col * 6.45)
    t = Inches(1.55 + row * 2.7)
    add_rect(sl, l, t, Inches(6.1), Inches(0.45), fill_rgb=color)
    add_text_box(sl, title, l+Inches(0.12), t+Inches(0.06),
                 Inches(5.9), Inches(0.36),
                 font_size=15, bold=True, color=C_WHITE)
    for j, b in enumerate(bullets):
        bt = t + Inches(0.5) + j * Inches(0.5)
        add_text_box(sl, "▪ " + b, l+Inches(0.15), bt,
                     Inches(5.8), Inches(0.45),
                     font_size=11, color=C_BLACK, wrap=True)


# ════════════════════════════════════════════════════
# SLIDE 7: 실험 결과 테이블
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LGRAY)
title_bar(sl, "04. 9개 모델 실험 결과", "Improved 방식 (BGE-m3-ko + BM25 + k=15 + 카테고리)")

headers = ["순위", "모델", "VRAM", "정확도", "속도", "처리시간"]
col_ws  = [Inches(0.6), Inches(3.4), Inches(1.1), Inches(1.4), Inches(1.2), Inches(1.5)]
col_ls  = [Inches(0.3)]
for w in col_ws[:-1]:
    col_ls.append(col_ls[-1] + w)

rows = [
    ("🥇 1", "gemma4:e4b (9.6GB)",      "9.6GB",  "88.47%", "1.58개/초", "35분"),
    ("🥈 2", "gemma3:4b (3.3GB)",        "3.3GB",  "84.48%", "1.77개/초", "31분"),
    ("🥉 3", "granite4.1:8b (5.3GB)",    "5.3GB",  "82.48%", "1.21개/초", "46분"),
    ("  4",  "gemma4:e2b (7.2GB)",       "7.2GB",  "80.93%", "1.93개/초", "29분"),
    ("  5",  "exaone3.5:7.8b (4.8GB)",   "4.8GB",  "79.16%", "1.23개/초", "45분"),
    ("  6",  "deepseek-r1:1.5b (1.1GB)", "1.1GB",  "78.05%", "1.49개/초", "37분"),
    ("  7",  "deepseek-r1:8b (5.2GB)",   "5.2GB",  "77.83%", "1.40개/초", "40분"),
    ("  8",  "exaone3.5:2.4b (1.6GB)",   "1.6GB",  "74.06%", "1.58개/초", "35분"),
    ("  9",  "granite4.1:3b (2.1GB)",    "2.1GB",  "67.85%", "1.53개/초", "36분"),
]

t_hdr = Inches(1.45)
for j, (hdr, cw, cl) in enumerate(zip(headers, col_ws, col_ls)):
    add_rect(sl, cl, t_hdr, cw, Inches(0.38), fill_rgb=C_NAVY)
    add_text_box(sl, hdr, cl+Inches(0.05), t_hdr+Inches(0.05),
                 cw-Inches(0.1), Inches(0.3),
                 font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

for r, row in enumerate(rows):
    row_t = Inches(1.83) + r * Inches(0.53)
    bg = RGBColor(0xEB, 0xF5, 0xFF) if r == 0 else (C_WHITE if r % 2 == 0 else C_LGRAY)
    for j, (val, cw, cl) in enumerate(zip(row, col_ws, col_ls)):
        add_rect(sl, cl, row_t, cw, Inches(0.5), fill_rgb=bg)
        fc = C_GREEN if j == 3 and r == 0 else (C_ACCENT if j == 3 else C_BLACK)
        bold = (r == 0 and j in [1, 3])
        add_text_box(sl, val.strip(), cl+Inches(0.05), row_t+Inches(0.08),
                     cw-Inches(0.1), Inches(0.38),
                     font_size=11, bold=bold, color=fc, align=PP_ALIGN.CENTER)

add_text_box(sl, "▲ Baseline 대비 +33.26%p 향상 (55.21% → 88.47%)   |   모델 크기 ≠ 성능 — 아키텍처가 중요",
             Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.38),
             font_size=12, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 8: Trade-Off 분석
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LGRAY)
title_bar(sl, "04. 속도 vs 정확도 트레이드오프", "MAUT 프레임워크 적용 (Keeney, 1975) [19]")

# 수식
add_rect(sl, Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.65), fill_rgb=C_NAVY)
add_text_box(sl, "Trade-Off Score  =  α × 정확도  +  (1−α) × 속도_점수",
             Inches(0.5), Inches(1.56), Inches(12.3), Inches(0.5),
             font_size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# α 시나리오 카드
scenarios = [
    ("α = 0.75", "정확도 최우선",  "🥇 gemma4:e4b\n🥈 gemma3:4b",    C_ORANGE,
     "정제 정확도가 비즈니스 핵심\n(법적 요건, 고가 부품 관리)"),
    ("α = 0.50", "균형",           "🥇 gemma3:4b\n🥈 gemma4:e4b",    C_GREEN,
     "정확도·처리속도 동등 중요\n(일반 업무 데이터 정제)"),
    ("α = 0.25", "속도 최우선",    "🥇 gemma4:e2b\n🥈 gemma3:4b",    C_ACCENT,
     "대량·실시간 처리 우선\n(일일 자동 배치 파이프라인)"),
]
for i, (alpha, label, winner, color, use) in enumerate(scenarios):
    l = Inches(0.4 + i * 4.3)
    t = Inches(2.35)
    add_rect(sl, l, t, Inches(4.0), Inches(0.5), fill_rgb=color)
    add_text_box(sl, f"{alpha}  ({label})", l+Inches(0.1), t+Inches(0.08),
                 Inches(3.8), Inches(0.36),
                 font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, l, t+Inches(0.5), Inches(4.0), Inches(3.55),
             fill_rgb=C_WHITE, line_rgb=color, line_width=1.5)
    add_text_box(sl, winner, l+Inches(0.1), t+Inches(0.65),
                 Inches(3.8), Inches(1.1),
                 font_size=20, bold=True, color=color, align=PP_ALIGN.CENTER, wrap=True)
    add_text_box(sl, "추천 사용 시나리오:", l+Inches(0.15), t+Inches(1.85),
                 Inches(3.7), Inches(0.3),
                 font_size=11, bold=True, color=C_DKGRAY)
    add_text_box(sl, use, l+Inches(0.15), t+Inches(2.15),
                 Inches(3.7), Inches(0.85),
                 font_size=12, color=C_BLACK, wrap=True)

add_text_box(sl, "* Belcak et al. (2025) [7]: 1.1GB deepseek-r1:1.5b → 78% — 소형 모델도 실무 적용 가능",
             Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.38),
             font_size=11, color=C_DKGRAY, italic=True)


# ════════════════════════════════════════════════════
# SLIDE 9: 오류 유형 분석 — 프레임워크
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LGRAY)
title_bar(sl, "05. LLM은 어디서 실수하는가 — 오류 유형 분석", "9개 모델 842건 오답 분류")

# 핵심 수치
key_data = [
    ("842건",   "전체 오답\n(9개 모델 합산)", C_NAVY),
    ("52건",    "최우수 모델 오답\n(gemma4:e4b, 88.47%)", C_ACCENT),
    ("46.2%",   "RAG 단계 실패\n(구조적 문제)", C_RED),
    ("44.2%",   "LLM 추론 실패\n(모델 능력 차이)", C_ORANGE),
]
for i, (num, lbl, color) in enumerate(key_data):
    l = Inches(0.4 + i * 3.15)
    key_number(sl, num, lbl, l, Inches(1.5), w=Inches(2.85), h=Inches(1.6), num_color=color)

# 2단계 실패 모델
add_text_box(sl, "오류 발생 위치: 2단계 구조",
             Inches(0.4), Inches(3.3), Inches(12.5), Inches(0.38),
             font_size=15, bold=True, color=C_NAVY)

# 단계 박스
pipeline = [
    ("더티 공구명 입력",  C_NAVY,   ""),
    ("RAG 검색 (15개 후보 생성)",  C_ACCENT, "←  단계 1 실패\n정답이 후보에 없음\n(~25건 고정)"),
    ("LLM 분류 (후보 중 선택)",    C_GREEN,  "←  단계 2 실패\n정답이 후보에 있지만 틀림\n(모델별 23~120건)"),
    ("표준명 출력",  C_NAVY,   ""),
]
for i, (name, col, note) in enumerate(pipeline):
    l = Inches(0.4 + i * 2.85)
    t = Inches(3.8)
    add_rect(sl, l, t, Inches(2.5), Inches(0.55), fill_rgb=col)
    add_text_box(sl, name, l+Inches(0.07), t+Inches(0.08),
                 Inches(2.36), Inches(0.42),
                 font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=True)
    if i < len(pipeline)-1:
        add_text_box(sl, "▶", l+Inches(2.5), t+Inches(0.1),
                     Inches(0.35), Inches(0.38),
                     font_size=16, bold=True, color=C_DKGRAY, align=PP_ALIGN.CENTER)
    if note:
        add_rect(sl, l, t+Inches(0.6), Inches(2.5), Inches(1.55),
                 fill_rgb=C_WHITE, line_rgb=col, line_width=1.5)
        add_text_box(sl, note, l+Inches(0.1), t+Inches(0.68),
                     Inches(2.3), Inches(1.4),
                     font_size=11, color=col, bold=True, wrap=True)

add_text_box(sl,
    "핵심 발견:  RAG 누락 ~25건은 모든 모델에서 동일 (구조 문제) / LLM 오선택은 모델에 따라 최대 5배 차이",
    Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.38),
    font_size=12, bold=True, color=C_RED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 10: 오류 세부 유형 + 모델 비교
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LGRAY)
title_bar(sl, "05. 오류 세부 유형 + 모델 간 비교", "gemma4:e4b 52건 상세 + 전체 모델 비교")

# 좌: 세부 유형 바
add_text_box(sl, "gemma4:e4b 오류 세부 유형 (52건)",
             Inches(0.3), Inches(1.5), Inches(6.0), Inches(0.38),
             font_size=14, bold=True, color=C_NAVY)

err_types = [
    ("RAG 검색 실패",        24, 46.2, C_RED),
    ("완전 다른 공구 선택",  12, 23.1, C_ORANGE),
    ("세트·단품 혼동",        7, 13.5, C_ACCENT),
    ("동력원 분류 오류",      4,  7.7, C_GREEN),
    ("본체·부속품 혼동",      4,  7.7, C_NAVY),
    ("드릴·드라이버 혼동",    1,  1.9, C_DKGRAY),
]

max_bar = Inches(4.5)
for i, (name, cnt, pct, color) in enumerate(err_types):
    t = Inches(2.0) + i * Inches(0.77)
    bar_w = max_bar * (pct / 50.0)
    add_text_box(sl, name, Inches(0.3), t, Inches(2.6), Inches(0.4),
                 font_size=11, color=C_BLACK)
    add_rect(sl, Inches(2.95), t+Inches(0.05), bar_w, Inches(0.3), fill_rgb=color)
    add_text_box(sl, f"{cnt}건 ({pct}%)", Inches(2.95)+bar_w+Inches(0.1), t,
                 Inches(1.0), Inches(0.38), font_size=11, bold=True, color=color)

# 우: 모델 비교
add_text_box(sl, "모델별 오류 비교 (RAG 누락 vs LLM 오선택)",
             Inches(7.0), Inches(1.5), Inches(5.9), Inches(0.38),
             font_size=14, bold=True, color=C_NAVY)

model_cmp = [
    ("gemma4:e4b",          24, 28,  88.47),
    ("gemma3:4b",           25, 45,  84.48),
    ("granite4.1:8b",       25, 54,  82.48),
    ("deepseek-r1:1.5b",    25, 74,  78.05),
    ("granite4.1:3b",       25, 120, 67.85),
]

hdrs2 = ["모델", "RAG 누락", "LLM 오선택", "정확도"]
cw2   = [Inches(2.1), Inches(1.0), Inches(1.1), Inches(0.95)]
cl2   = [Inches(7.1)]
for w in cw2[:-1]: cl2.append(cl2[-1]+w)

th = Inches(1.98)
for j, (h, cw, cl) in enumerate(zip(hdrs2, cw2, cl2)):
    bg = C_NAVY if j == 0 else (C_RED if j == 1 else (C_ORANGE if j == 2 else C_GREEN))
    add_rect(sl, cl, th, cw, Inches(0.35), fill_rgb=bg)
    add_text_box(sl, h, cl+Inches(0.04), th+Inches(0.04),
                 cw-Inches(0.08), Inches(0.28),
                 font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

for r, (mdl, rag, llm, acc) in enumerate(model_cmp):
    rt = Inches(2.33) + r * Inches(0.78)
    bg = RGBColor(0xEB, 0xF5, 0xFF) if r == 0 else (C_WHITE if r % 2 == 0 else C_LGRAY)
    for j, (val, cw, cl) in enumerate(zip([mdl, str(rag), str(llm), f"{acc}%"],
                                           cw2, cl2)):
        add_rect(sl, cl, rt, cw, Inches(0.7), fill_rgb=bg)
        c = C_BLACK if j == 0 else (C_RED if j == 1 else (C_ORANGE if j == 2 else C_GREEN))
        add_text_box(sl, val, cl+Inches(0.04), rt+Inches(0.15),
                     cw-Inches(0.08), Inches(0.42),
                     font_size=11, bold=(r == 0), color=c, align=PP_ALIGN.CENTER)

add_text_box(sl, "→ RAG 개선(~25건 고정 문제) + LLM 프롬프트 특화(세트/단품, 동력원) = 이중 전략 필요",
             Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.38),
             font_size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 11: 자주 틀리는 패턴 + 입력 유형
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LGRAY)
title_bar(sl, "05. 자주 틀리는 패턴 + 오류 기여 입력 유형", "취약점 심층 분석")

# 좌: 반복 오류 사례
add_text_box(sl, "반복 오류 사례 (9개 모델 공통)",
             Inches(0.3), Inches(1.5), Inches(6.5), Inches(0.38),
             font_size=14, bold=True, color=C_NAVY)

patterns = [
    ("몽키스패너 → 스패너",         "GT: 몽키렌치 / 동의어 혼동"),
    ("L렌치세트 → T렌치세트",       "GT: 육각렌치세트 / L형·T형 혼동"),
    ("노기스 → (RAG 누락)",         "GT: 버니어캘리퍼스 / 일본어 잔재"),
    ("기리세트 → 세트",             "GT: 드릴비트세트 / 일본어+세트"),
    ("강력절단기 → 고속절단기",      "GT: 볼트컷터 / 절단기 계열 혼동"),
    ("글루건&글루스틱 → 글루스틱",   "GT: 글루건 / 복합명 부속품 선택"),
    ("bosch2(무선) → 전동드릴",     "GT: 충전드릴 / 무선=충전 미인식"),
]
for i, (pat, note) in enumerate(patterns):
    t = Inches(2.0) + i * Inches(0.67)
    add_rect(sl, Inches(0.3), t, Inches(6.2), Inches(0.58),
             fill_rgb=C_WHITE, line_rgb=C_DKGRAY, line_width=0.5)
    add_text_box(sl, pat,  Inches(0.4),  t+Inches(0.06), Inches(3.0), Inches(0.28),
                 font_size=11, bold=True, color=C_RED)
    add_text_box(sl, note, Inches(3.45), t+Inches(0.06), Inches(3.0), Inches(0.28),
                 font_size=10, color=C_DKGRAY)

# 우: 입력 유형별 오류 기여도
add_text_box(sl, "오답 행 입력 유형 분포 (gemma4:e4b, 52건)",
             Inches(7.0), Inches(1.5), Inches(5.9), Inches(0.38),
             font_size=14, bold=True, color=C_NAVY)

input_types = [
    ("단순 명사형",   17, 32.7, C_RED),
    ("세트·복합",     14, 26.9, C_ORANGE),
    ("브랜드명 혼입", 10, 19.2, C_ACCENT),
    ("일본어 잔재",   10, 19.2, C_NAVY),
    ("모델번호 포함",  8, 15.4, C_DKGRAY),
    ("규격·속성",      6, 11.5, C_GREEN),
    ("동력원 명시",    4,  7.7, C_GREEN),
]
max_bar2 = Inches(4.0)
for i, (name, cnt, pct, color) in enumerate(input_types):
    t = Inches(2.0) + i * Inches(0.67)
    bar_w = max_bar2 * (pct / 35.0)
    add_text_box(sl, name, Inches(7.0), t, Inches(1.9), Inches(0.38),
                 font_size=11, color=C_BLACK)
    add_rect(sl, Inches(9.0), t+Inches(0.05), bar_w, Inches(0.28), fill_rgb=color)
    add_text_box(sl, f"{pct}%", Inches(9.0)+bar_w+Inches(0.08), t,
                 Inches(0.7), Inches(0.35), font_size=11, bold=True, color=color)

add_text_box(sl, "* 단순 명사형 오답 17건 — 입력이 단순해도 오탈자·약어로 RAG 실패 또는 LLM 혼동 발생",
             Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.38),
             font_size=11, color=C_DKGRAY, italic=True)


# ════════════════════════════════════════════════════
# SLIDE 12: 구현 로드맵
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LGRAY)
title_bar(sl, "06. 구현 로드맵", "오류 유형 분석 기반 4단계 도입 계획")

steps2 = [
    ("Week 1–2",
     "정답지(Ground Truth) 구축",
     ["도메인 전문가 + 카탈로그 참조",
      "500~1,000개 엔트리 목표",
      "일본어 잔재·오탈자 사전 병행 구축",
      "→ RAG 검색 실패 ~25건 구조 문제 해결 준비"],
     C_ACCENT),
    ("Week 3",
     "기술 스택 설치",
     ["Ollama, sentence-transformers",
      "faiss-cpu, rank-bm25",
      "Python 환경 구성",
      "→ 로컬 실행 환경 완성"],
     C_GREEN),
    ("Week 4",
     "시범 운영 + 오류 유형 진단",
     ["1,000건 샘플 정확도 측정",
      "오류 유형 분류: RAG 실패 vs LLM 실패",
      "본 연구 프레임워크 적용",
      "→ 취약 패턴 파악"],
     C_ORANGE),
    ("Month 2+",
     "유형별 개선 + 전체 배치",
     ["RAG 실패: 전처리 사전 확충",
      "LLM 실패: 세트/단품·동력원 특화 프롬프트",
      "Human-in-the-Loop: 낮은 신뢰도 항목 검수 [3]",
      "→ 신규 데이터 자동 파이프라인 연동"],
     C_RED),
]

colors = [C_ACCENT, C_GREEN, C_ORANGE, C_RED]
for i, (period, title, bullets, color) in enumerate(steps2):
    l = Inches(0.3 + i * 3.25)
    add_rect(sl, l, Inches(1.5), Inches(3.0), Inches(0.45), fill_rgb=color)
    add_text_box(sl, period, l+Inches(0.08), Inches(1.55),
                 Inches(2.85), Inches(0.35),
                 font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, l, Inches(1.95), Inches(3.0), Inches(0.42),
             fill_rgb=RGBColor(0xE8, 0xEE, 0xF5))
    add_text_box(sl, title, l+Inches(0.08), Inches(1.97),
                 Inches(2.85), Inches(0.38),
                 font_size=12, bold=True, color=C_NAVY)
    for j, b in enumerate(bullets):
        bt = Inches(2.44) + j * Inches(0.98)
        bullet_sym = "▪ " if not b.startswith("→") else ""
        fc = C_BLACK if not b.startswith("→") else color
        add_text_box(sl, bullet_sym + b, l+Inches(0.12), bt,
                     Inches(2.8), Inches(0.88),
                     font_size=10, color=fc, bold=b.startswith("→"), wrap=True)


# ════════════════════════════════════════════════════
# SLIDE 13: 기대 효과 + 한계
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LGRAY)
title_bar(sl, "06. 기대 효과 및 잔여 한계", "도입 전·후 비교 + 개선 과제")

# 좌: 기대 효과
add_text_box(sl, "기대 효과",
             Inches(0.3), Inches(1.5), Inches(5.8), Inches(0.38),
             font_size=15, bold=True, color=C_GREEN)

effects = [
    ("처리 속도", "3~5일/3,000건", "35분/3,000건", "×100+ 향상"),
    ("인력 투입", "담당자 상시", "최초 정답지만", "고정비용화"),
    ("일관성",   "담당자 편차", "100% 동일 규칙", "품질 표준화"),
    ("정확도",   "숙련자 ~95%", "RAG: 74~88%", "자동화 절충"),
]
eff_hdrs = ["지표", "현재(수작업)", "도입 후", "효과"]
ew = [Inches(1.2), Inches(1.5), Inches(1.5), Inches(1.3)]
el = [Inches(0.3)]
for w in ew[:-1]: el.append(el[-1]+w)

for j, (h, w, l) in enumerate(zip(eff_hdrs, ew, el)):
    add_rect(sl, l, Inches(1.95), w, Inches(0.35), fill_rgb=C_GREEN)
    add_text_box(sl, h, l+Inches(0.04), Inches(1.98), w-Inches(0.08), Inches(0.28),
                 font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

for r, row in enumerate(effects):
    rt = Inches(2.3) + r * Inches(0.72)
    bg = C_WHITE if r % 2 == 0 else C_LGRAY
    for j, (val, w, l) in enumerate(zip(row, ew, el)):
        add_rect(sl, l, rt, w, Inches(0.65), fill_rgb=bg)
        c = C_BLACK if j < 2 else (C_GREEN if j == 2 else C_ACCENT)
        add_text_box(sl, val, l+Inches(0.04), rt+Inches(0.1),
                     w-Inches(0.08), Inches(0.48),
                     font_size=10, color=c, bold=(j >= 2), align=PP_ALIGN.CENTER, wrap=True)

# 우: 잔여 한계 + 개선 방향
add_text_box(sl, "잔여 한계 및 개선 방향",
             Inches(6.4), Inches(1.5), Inches(6.5), Inches(0.38),
             font_size=15, bold=True, color=C_RED)

limits = [
    ("RAG 구조 한계 (~25건)",
     "일본어 잔재·모델번호 입력 → RAG 검색 실패\n→ 전처리 레이어 추가로 개선 가능", C_RED),
    ("세트·단품 혼동 (13.5%)",
     "복합 공구명에서 주공구 판단 실패\n→ 프롬프트 엔지니어링으로 개선 가능", C_ORANGE),
    ("동력원 분류 오류 (7.7%)",
     "충전식=무선 등 동의어 미인식\n→ 동의어 사전 추가로 해결", C_ACCENT),
    ("대형 모델 실행 불가",
     "RTX 5080(17.1GB): 30B+ 모델 0.03개/초\n→ 클라우드 또는 대형 GPU 서버 필요", C_NAVY),
]
for i, (title, desc, color) in enumerate(limits):
    t = Inches(1.95) + i * Inches(1.3)
    add_rect(sl, Inches(6.4), t, Inches(0.1), Inches(1.1), fill_rgb=color)
    add_text_box(sl, title, Inches(6.6), t+Inches(0.06),
                 Inches(6.2), Inches(0.3),
                 font_size=12, bold=True, color=color)
    add_text_box(sl, desc,  Inches(6.6), t+Inches(0.38),
                 Inches(6.2), Inches(0.65),
                 font_size=10, color=C_BLACK, wrap=True)


# ════════════════════════════════════════════════════
# SLIDE 14: 결론
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_NAVY)

# 상단 바
add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_rgb=C_ACCENT)

add_text_box(sl, "결론 및 제언",
             Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.65),
             font_size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

conclusions = [
    ("① 비정형 데이터 정제는 LLM+RAG로 해결 가능하다",
     "Baseline 55% → Improved 88.47% (+33%p)\nKim et al.[18], Zhang & Huang[12], Glock & Korom[10]",
     C_ACCENT),
    ("② 오류 유형 분석이 개선의 핵심이다",
     "RAG 누락 46% (구조 문제) + LLM 오선택 44% (추론 문제)\n이중 전략: RAG 개선 + LLM 프롬프트 특화",
     C_ORANGE),
    ("③ 로컬 LLM은 클라우드 API의 실질적 대안이다",
     "보안·비용·안정성 모두 우위\nBelcak[7], WINGPT[2], Amugongo[5]",
     C_GREEN),
]

for i, (title, body, color) in enumerate(conclusions):
    t = Inches(1.1) + i * Inches(1.85)
    add_rect(sl, Inches(0.4), t, Inches(0.08), Inches(1.55), fill_rgb=color)
    add_text_box(sl, title, Inches(0.6), t+Inches(0.08),
                 Inches(12.1), Inches(0.5),
                 font_size=16, bold=True, color=C_WHITE)
    add_text_box(sl, body,  Inches(0.6), t+Inches(0.6),
                 Inches(12.1), Inches(0.85),
                 font_size=13, color=RGBColor(0xAE, 0xC6, 0xCF), wrap=True)

# 향후 연구
add_rect(sl, Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.55),
         fill_rgb=RGBColor(0x0D, 0x1B, 0x36))
add_text_box(sl, "향후 연구:  일본어 전처리 레이어  |  세트/단품 특화 프롬프트  "
                 "|  정답지 3,352건 전체 확장  |  파인튜닝 비교  |  Human-in-the-Loop 통합",
             Inches(0.5), Inches(6.68), Inches(12.3), Inches(0.45),
             font_size=11, color=RGBColor(0xAE, 0xC6, 0xCF), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════
OUT = "/mnt/c/Users/neohc/Desktop/ClaudeCode/_workspace5/발표초안_v1.pptx"
prs.save(OUT)
print(f"저장 완료: {OUT}")
print(f"슬라이드 수: {len(prs.slides)}장")
